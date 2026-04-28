"""PPTX builder for briefing slides."""

from __future__ import annotations

from io import BytesIO
from pathlib import Path
from urllib.request import urlopen

from src.schemas import ArticleSummary, SlidePlan


class PPTXBuildError(RuntimeError):
    """Raised when PPTX output cannot be created."""


FONT_NAME = "나눔스퀘어 Bold"
COLOR_BLACK = "000000"
COLOR_CONTENT_BG = "F2F2F2"
COLOR_BORDER = "AFABAB"
COLOR_META = "77787B"
COLOR_BLUE = "2212FF"


def build_briefing_pptx(
    summaries: tuple[ArticleSummary, ...],
    slide_plans: tuple[SlidePlan, ...],
    output_path: Path,
) -> Path:
    """Build a PPTX briefing deck and return the output path."""
    _validate_output_path(output_path)
    if not summaries:
        raise PPTXBuildError("at least one article summary is required")
    for summary in summaries:
        _validate_summary(summary)
    try:
        from pptx import Presentation
        from pptx.dml.color import RGBColor
        from pptx.enum.shapes import MSO_SHAPE
        from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
        from pptx.util import Inches, Pt
    except ImportError as exc:
        raise PPTXBuildError("python-pptx package is not installed") from exc

    output_path.parent.mkdir(parents=True, exist_ok=True)
    rgb = _rgb_factory(RGBColor)
    presentation = _load_presentation(Presentation, Inches)
    plan_by_article = {plan.article_id: plan for plan in slide_plans}

    for index, summary in enumerate(summaries, start=1):
        slide = presentation.slides.add_slide(_blank_layout(presentation))
        _add_header(slide, summary, index, Inches, Pt, rgb, MSO_SHAPE, MSO_ANCHOR, PP_ALIGN)
        _add_content_box(slide, Inches, rgb, MSO_SHAPE)
        _add_issue_content(
            slide,
            summary,
            plan_by_article.get(summary.article_id),
            Inches,
            Pt,
            rgb,
            PP_ALIGN,
        )
        _add_footer(slide, summary, Inches, Pt, rgb, PP_ALIGN)

    presentation.save(output_path)
    return output_path


def _load_presentation(presentation_factory: object, inches: object) -> object:
    template_path = _find_template_pptx()
    if template_path is None:
        presentation = presentation_factory()
    else:
        presentation = presentation_factory(str(template_path))
        _clear_existing_slides(presentation)
    presentation.slide_width = inches(13.333)
    presentation.slide_height = inches(7.5)
    return presentation


def _find_template_pptx() -> Path | None:
    templates_dir = Path("templates")
    if not templates_dir.exists():
        return None
    candidates = sorted(
        path
        for pattern in ("*.pptx", "*.potx")
        for path in templates_dir.glob(pattern)
        if not path.name.startswith("~$")
    )
    return candidates[0] if candidates else None


def _rgb_factory(rgb_class: object) -> object:
    def make(hex_color: str) -> object:
        return rgb_class(
            int(hex_color[0:2], 16),
            int(hex_color[2:4], 16),
            int(hex_color[4:6], 16),
        )

    return make


def _clear_existing_slides(presentation: object) -> None:
    slide_id_list = presentation.slides._sldIdLst
    for slide_id in list(slide_id_list):
        presentation.part.drop_rel(slide_id.rId)
        slide_id_list.remove(slide_id)


def _blank_layout(presentation: object) -> object:
    for layout in presentation.slide_layouts:
        if getattr(layout, "name", "").lower() == "blank":
            return layout
    return presentation.slide_layouts[-1]


def _add_header(
    slide: object,
    summary: ArticleSummary,
    index: int,
    inches: object,
    pt: object,
    rgb: object,
    shape_enum: object,
    anchor_enum: object,
    align_enum: object,
) -> None:
    header = slide.shapes.add_shape(
        shape_enum.RECTANGLE,
        inches(0.3),
        inches(0.25),
        inches(12.73),
        inches(0.95),
    )
    _style_rect(header, rgb, COLOR_CONTENT_BG, COLOR_BORDER)

    title = _add_textbox(
        slide,
        f"{index:02d}. {summary.title}",
        inches(0.55),
        inches(0.36),
        inches(12.2),
        inches(0.45),
        pt(20),
        rgb(COLOR_BLACK),
        align_enum.CENTER,
        anchor_enum.MIDDLE,
    )
    title.margin_left = inches(0.06)
    title.margin_right = inches(0.06)

    _add_textbox(
        slide,
        f"{summary.source} | {summary.pdf_source}",
        inches(5.4),
        inches(0.88),
        inches(7.25),
        inches(0.27),
        pt(13),
        rgb(COLOR_META),
        align_enum.RIGHT,
        anchor_enum.BOTTOM,
    )

    _add_textbox(
        slide,
        _headline(summary.summary),
        inches(0.55),
        inches(1.28),
        inches(12.2),
        inches(0.35),
        pt(14),
        rgb(COLOR_BLUE),
        align_enum.LEFT,
        anchor_enum.MIDDLE,
    )


def _add_content_box(slide: object, inches: object, rgb: object, shape_enum: object) -> None:
    content_box = slide.shapes.add_shape(
        shape_enum.RECTANGLE,
        inches(0.35),
        inches(1.75),
        inches(12.6),
        inches(5.38),
    )
    _style_rect(content_box, rgb, COLOR_CONTENT_BG, COLOR_BORDER)


def _add_issue_content(
    slide: object,
    summary: ArticleSummary,
    plan: SlidePlan | None,
    inches: object,
    pt: object,
    rgb: object,
    align_enum: object,
) -> None:
    left_text = (
        "1. 배경 / 맥락\n"
        f"{summary.pdf_summary}\n\n"
        "2. 핵심 팩트\n"
        f"{summary.summary}"
    )
    _add_textbox(
        slide,
        left_text,
        inches(0.55),
        inches(1.95),
        inches(6.3),
        inches(5.05),
        pt(10.2),
        rgb(COLOR_BLACK),
        align_enum.LEFT,
        None,
    )

    image_added = _try_add_article_image(
        slide,
        summary,
        inches(7.0),
        inches(1.95),
        inches(5.75),
        inches(2.35),
    )
    right_top = inches(4.45) if image_added else inches(1.95)
    right_height = inches(2.55) if image_added else inches(5.05)
    right_text = (
        "3. [당사 함의] So What\n"
        f"{summary.conclusion}\n\n"
        "4. 당사 Action Items\n"
        f"{_action_items(summary)}"
    )
    if not image_added and plan is not None and plan.visual_type:
        right_text = f"{right_text}\n\nVisual: {plan.visual_type}"
    _add_textbox(
        slide,
        right_text,
        inches(7.0),
        right_top,
        inches(5.75),
        right_height,
        pt(10.2),
        rgb(COLOR_BLACK),
        align_enum.LEFT,
        None,
    )


def _add_footer(
    slide: object,
    summary: ArticleSummary,
    inches: object,
    pt: object,
    rgb: object,
    align_enum: object,
) -> None:
    _add_textbox(
        slide,
        f"기사 원본URL: {summary.url}",
        inches(0.45),
        inches(7.14),
        inches(8.2),
        inches(0.18),
        pt(7.2),
        rgb(COLOR_META),
        align_enum.LEFT,
        None,
    )
    _add_textbox(
        slide,
        (
            "삼우씨엠(SAMOO CM) 전략사업그룹 | "
            "2026-W18 Weekly Construction Briefing | 보고용 브리핑 자료"
        ),
        inches(5.7),
        inches(7.32),
        inches(7.15),
        inches(0.18),
        pt(7.2),
        rgb(COLOR_META),
        align_enum.RIGHT,
        None,
    )


def _add_textbox(
    slide: object,
    text: str,
    left: object,
    top: object,
    width: object,
    height: object,
    font_size: object,
    color: object,
    alignment: object | None,
    vertical_anchor: object | None,
) -> object:
    box = slide.shapes.add_textbox(left, top, width, height)
    frame = box.text_frame
    frame.clear()
    frame.word_wrap = True
    frame.margin_left = 0
    frame.margin_right = 0
    frame.margin_top = 0
    frame.margin_bottom = 0
    if vertical_anchor is not None:
        frame.vertical_anchor = vertical_anchor
    frame.text = text
    for paragraph in frame.paragraphs:
        if alignment is not None:
            paragraph.alignment = alignment
        paragraph.font.name = FONT_NAME
        paragraph.font.size = font_size
        paragraph.font.bold = False
        paragraph.font.color.rgb = color
        for run in paragraph.runs:
            run.font.name = FONT_NAME
            run.font.size = font_size
            run.font.bold = False
            run.font.color.rgb = color
    return frame


def _style_rect(shape: object, rgb: object, fill_color: str, line_color: str) -> None:
    shape.fill.solid()
    shape.fill.fore_color.rgb = rgb(fill_color)
    shape.line.color.rgb = rgb(line_color)
    shape.line.width = 3175


def _headline(summary_text: str) -> str:
    first_line = summary_text.splitlines()[0] if summary_text else ""
    return first_line.replace("핵심:", "", 1).strip() or "기사 핵심 요약"


def _action_items(summary: ArticleSummary) -> str:
    topic = summary.topic
    if "원가" in topic or "자재" in topic or "PF" in topic:
        return (
            "- 진행 현장 원가·공기 영향도 재점검\n"
            "- 계약상 물가변동·공기연장 조항 확인\n"
            "- 발주처 협의 자료와 리스크 시나리오 준비"
        )
    if "수주" in topic or "입찰" in topic:
        return (
            "- 입찰 기준 변경 가능성 모니터링\n"
            "- 견적 산출 근거와 의사결정 기록 정비\n"
            "- 고위험 발주 건은 사전 법무·원가 검토"
        )
    if "기술" in topic or "BIM" in topic or "스마트" in topic:
        return (
            "- 적용 가능한 파일럿 현장 선정\n"
            "- 데이터·표준·협력사 역량 점검\n"
            "- 생산성 개선 효과를 비용·일정 지표로 추적"
        )
    if "시장" in topic or "부동산" in topic or "미분양" in topic:
        return (
            "- 지역별 수요·가격 가정 업데이트\n"
            "- 보유자산과 미분양 현금흐름 점검\n"
            "- 투자·분양·매각 시나리오 재검토"
        )
    return (
        "- 관련 부서 영향도 확인\n"
        "- 대응 과제와 책임자 지정\n"
        "- 다음 주 브리핑에서 정책·시장 변화 추적"
    )


def _validate_output_path(output_path: Path) -> None:
    normalized = output_path.as_posix()
    if "/data/output/" not in f"/{normalized}":
        raise PPTXBuildError("PPTX output must be written under data/output")


def _validate_summary(summary: ArticleSummary) -> None:
    if not summary.source or not summary.url:
        raise PPTXBuildError(
            f"slide metadata missing for article {summary.article_id}: source and url required"
        )
    if (
        not summary.pdf_source
        or not summary.topic
        or not summary.pdf_summary
        or not summary.conclusion
    ):
        raise PPTXBuildError(
            f"slide metadata missing for article {summary.article_id}: "
            "pdf_source, topic, pdf_summary, and conclusion required"
        )


def _try_add_article_image(
    slide: object,
    summary: ArticleSummary,
    left: object,
    top: object,
    max_width: object,
    max_height: object,
) -> bool:
    if not summary.image_url:
        return False
    try:
        image_data = _image_bytes(summary.image_url)
        picture = slide.shapes.add_picture(image_data, left, top, width=max_width)
        if picture.height > max_height:
            ratio = max_height / picture.height
            picture.height = max_height
            picture.width = int(picture.width * ratio)
        if picture.width > max_width:
            ratio = max_width / picture.width
            picture.width = max_width
            picture.height = int(picture.height * ratio)
        picture.left = int(left + (max_width - picture.width) / 2)
        picture.top = int(top + (max_height - picture.height) / 2)
    except Exception:
        return False
    return True


def _image_bytes(image_url: str) -> BytesIO | str:
    if image_url.startswith(("http://", "https://")):
        with urlopen(image_url, timeout=10) as response:
            return BytesIO(response.read())
    path = Path(image_url)
    if not path.exists():
        raise FileNotFoundError(image_url)
    return str(path)
